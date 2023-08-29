from pptx import Presentation
from config.config import endpoint, headers
import requests

def taranslate_presentation_to_hindi(input_presentation_file, output_ppt_file):

    ppt = Presentation(input_presentation_file)

    # Iterate through slide text and translate
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # API request to translate
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        english_text = run.text
                        if english_text == "" or english_text == "\n" or english_text.strip() == "":
                            translated_text = english_text
                        else:
                            data = {
                                "model": "gpt-4",
                                "messages": [{
                                    "role": "user",
                                    "content": f"translate content from English to hindi: {english_text}"
                                }],
                            }
                            response = requests.post(
                                endpoint, json=data, headers=headers)
                            result = response.json()
                            translated_text = result['choices'][0]['message']['content']
                        # Update the shape with translated text
                        run.text = translated_text  

    # Save the translated PowerPoint
    translated_pptx_file = output_ppt_file
    ppt.save(translated_pptx_file)